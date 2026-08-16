#!/usr/bin/env python3
"""PowerPoint から、文字・表・ノートを markdown の下書きとして取り出す。

pandoc も .pptx を読めるが、**発表者ノートを落とし**、図形の文字を z-order 順に
平坦化する。本スクリプトは python-pptx で読み、

* スライドのタイトルを見出しにする
* 図形の文字を**読み順（上から下、左から右）**で並べる
* 表を markdown の表にする（結合セルは原点にだけ文字を置く）
* 発表者ノートをコメントとして残す
* 画像を書き出す
* **図形とコネクタで描かれたスライドには警告を入れる**
* 入力／処理／出力の3欄があるスライドは ``::: {.ipo}`` の下書きにする

図はここでは取り出せない
------------------------
枠・矢印・配置は markdown に写せない。取り出せるのは**文字だけ**である。
図そのものは提出用 PDF から ``pdfsvg.py`` で回収するか、mermaid ／ ``.ipo`` に
作り直すこと。本スクリプトはその判断ができるよう、図形とコネクタの数を
スライドごとに書き添える。

必要なもの
----------
python-pptx（``pip install python-pptx``）。

使い方
------
    python migration/pptxtext.py 資料.pptx -o draft.md
    python migration/pptxtext.py 資料.pptx -o draft.md --media docs/diagrams/ppt
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("エラー: python-pptx がありません。`pip install python-pptx` を実行してください。",
          file=sys.stderr)
    raise SystemExit(1)

# 行頭の箇条書き記号（PowerPoint の本文にそのまま入っていることが多い）
BULLET = re.compile(r"^[・･•●○■□\-\*–]\s*")

# 既に番号付きリストの形になっている行（「1. …」「(1) …」「① …」）
NUMBERED = re.compile(r"^(?:\d+[.)]|[(（]\d+[)）]|[①-⑳])\s+")

# IPO の欄名
IPO_SECTIONS = ("入力", "処理", "出力")

# 図が主体とみなす図形数
DIAGRAM_SHAPES = 6

# 読み順にまとめるときの縦方向の許容（EMU。約 0.2 インチ）
ROW_TOLERANCE = 182880


def _iter_shapes(shapes):
    """グループの中まで辿って、すべての図形を返す。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _reading_order(shapes) -> list:
    """上から下、左から右に並べ替える。"""
    def key(shape):
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
        return (round(top / ROW_TOLERANCE), left)
    return sorted(shapes, key=key)


def _lines_of(shape) -> list:
    """図形の中の文字を、箇条書きの階層つきで取り出す。"""
    if not getattr(shape, "has_text_frame", False):
        return []
    out = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            continue
        out.append(("  " * max(0, paragraph.level), BULLET.sub("", text)))
    return out


def _as_markdown(lines, force_list: bool = False) -> list:
    """(インデント, 文字) の並びを markdown の段落・箇条書きにする。"""
    out = []
    for indent, text in lines:
        if NUMBERED.match(text):
            # 「1. 入力値を検証」は既に番号付きリストの形なので、記号を重ねない
            out.append(f"{indent}{text}")
        elif force_list or indent:
            out.append(f"{indent}- {text}")
        else:
            out.append(text)
    return out


def _table_markdown(table) -> list:
    """PowerPoint の表を markdown のパイプ表にする。"""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            # 結合されたセルは原点にだけ文字を残す
            text = "" if getattr(cell, "is_spanned", False) else cell.text.strip()
            cells.append(text.replace("\n", "<br>").replace("|", "\\|"))
        rows.append(cells)
    if not rows:
        return []

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |",
           "|" + "|".join("---" for _ in range(width)) + "|"]
    out.extend("| " + " | ".join(r) + " |" for r in rows[1:])
    return out


def _find_ipo(shapes) -> dict:
    """入力／処理／出力の欄を持つスライドなら、その中身を返す。"""
    found = {}
    for shape in shapes:
        lines = _lines_of(shape)
        if not lines:
            continue
        head = lines[0][1].strip().rstrip(":：")
        if head in IPO_SECTIONS and head not in found:
            found[head] = lines[1:]
    return found if len(found) == len(IPO_SECTIONS) else {}


def convert(path: Path, media_dir: Path | None, image_prefix: str):
    prs = Presentation(str(path))
    out, summary = [], []

    for number, slide in enumerate(prs.slides, 1):
        shapes = list(_iter_shapes(slide.shapes))
        ordered = _reading_order(shapes)

        # python-pptx はアクセスのたびに別のラッパを返すため、`is` ではなく id で照合する
        title, title_id = "", None
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip().replace("\n", " ")
            title_id = slide.shapes.title.shape_id

        pictures = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        tables = [s for s in shapes if getattr(s, "has_table", False)]
        connectors = [s for s in shapes
                      if s.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM)
                      or s.__class__.__name__ == "Connector"]
        texted = [s for s in ordered if _lines_of(s)]
        is_diagram = len(connectors) > 0 or len(texted) >= DIAGRAM_SHAPES

        out.append(f"## スライド {number}　{title}" if title else f"## スライド {number}")
        out.append("")

        if is_diagram:
            out.append(f"<!-- 図が主体のスライド（文字のある図形 {len(texted)} 個 / "
                       f"コネクタ {len(connectors)} 本）。")
            out.append("     枠と矢印はここに写っていない。提出用 PDF の同じページから")
            out.append("     pdfsvg.py で SVG を取り出すか、mermaid / .ipo に作り直すこと。 -->")
            out.append("")

        ipo = _find_ipo(ordered)
        if ipo:
            caption = title or "（タイトルを書く）"
            out.append(f'::: {{.ipo module="（機能名を書く）" '
                       f'caption="{caption}" label="tbl-ppt-{number:03d}"}}')
            out.append("")
            out.append(f"## {title or '（処理名を書く）'}")
            out.append("")
            for section in IPO_SECTIONS:
                out.append(f"### {section}")
                out.append("")
                out.extend(_as_markdown(ipo[section], force_list=True) or ["（内容を書く）"])
                out.append("")
            out.append(":::")
            out.append("")

            # IPO 欄の外にある表は、どの欄に入るか機械では決められない。
            # 落とさずに後ろへ出し、置き場所を人が決める。
            for shape in ordered:
                if getattr(shape, "has_table", False):
                    out.append("<!-- 下の表は IPO のどの欄に入るか判断できなかった。"
                               "入力／処理／出力のいずれかへ移すこと。 -->")
                    out.append("")
                    out.extend(_table_markdown(shape.table))
                    out.append("")
            summary.append((number, title, "IPO の下書きにした", len(texted), len(connectors)))
        else:
            body = []
            for shape in ordered:
                if title_id is not None and shape.shape_id == title_id:
                    continue
                if getattr(shape, "has_table", False):
                    body.extend(_table_markdown(shape.table))
                    body.append("")
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                lines = _lines_of(shape)
                if lines:
                    # PowerPoint の本文は既定で箇条書きである。2行以上あるものは
                    # 箇条書きとみなす（1行だけの図形は段落のまま）。
                    body.extend(_as_markdown(lines, force_list=len(lines) >= 2))
                    body.append("")
            out.extend(body if body else ["（文字のないスライド）", ""])
            summary.append((number, title, "図が主体" if is_diagram else "文字と表",
                            len(texted), len(connectors)))

        # ---- 画像
        for index, picture in enumerate(pictures, 1):
            name = f"slide{number:03d}-{index}"
            if media_dir is not None:
                image = picture.image
                target = media_dir / f"{name}.{image.ext}"
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_bytes(image.blob)
                out.append(f"![（キャプションを書く）]({image_prefix}/{target.name})"
                           f"{{#fig-{name} width=80%}}")
            else:
                out.append(f"<!-- 画像 {index} 枚目。--media を付けると書き出す -->")
            out.append("")

        # ---- 発表者ノート
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append("<!-- 発表者ノート:")
                out.extend(f"     {line}" for line in notes.splitlines())
                out.append("-->")
                out.append("")

    return "\n".join(out).rstrip() + "\n", summary


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(
        description="PowerPoint から文字・表・ノートを markdown の下書きとして取り出す。",
        epilog="図そのものは pdfsvg.py で PDF から回収する。詳しくは migration/README.md を参照。")
    parser.add_argument("source", help="入力の .pptx")
    parser.add_argument("-o", "--out", help="出力先の markdown（既定: 標準出力）")
    parser.add_argument("--media", metavar="DIR", help="画像の書き出し先（例: docs/diagrams/ppt）")
    parser.add_argument("--image-prefix", default="/diagrams",
                        help="markdown に書く画像パスの前置き（既定: /diagrams）")
    args = parser.parse_args(argv)

    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    source = Path(args.source)
    if not source.exists():
        print(f"エラー: {source} がありません。", file=sys.stderr)
        return 1

    media_dir = Path(args.media) if args.media else None
    text, summary = convert(source, media_dir, args.image_prefix.rstrip("/"))

    if args.out:
        Path(args.out).write_text(text, encoding="utf-8")
    else:
        sys.stdout.write(text)

    print(f"\nスライド {len(summary)} 枚を処理しました。", file=sys.stderr)
    print("  枚  種別                図形  コネクタ  タイトル", file=sys.stderr)
    for number, title, kind, shapes, connectors in summary:
        print(f"  {number:>2}  {kind:<18} {shapes:>4} {connectors:>9}  {title[:30]}",
              file=sys.stderr)

    diagrams = [s for s in summary if s[2] == "図が主体"]
    if diagrams:
        pages = ",".join(str(s[0]) for s in diagrams)
        print(f"\n図が主体のスライドが {len(diagrams)} 枚あります。", file=sys.stderr)
        print("提出用 PDF の同じページから SVG を取り出す場合:", file=sys.stderr)
        print(f"  python migration/pdfsvg.py 提出版.pdf -o docs/diagrams --pages {pages}",
              file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
